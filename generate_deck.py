from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x2B, 0x55)   # slide backgrounds / headings
TEAL      = RGBColor(0x00, 0x8C, 0x8C)   # accent / highlights
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF4, 0xF8, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
AMBER     = RGBColor(0xF3, 0x9C, 0x12)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def bullet_box(slide, items, l, t, w, h, font_size=15, color=DARK_TEXT, indent=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        prefix = "  •  " if indent else ""
        run.text = prefix + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)

add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 5.8, 13.33, 1.7, TEAL)

add_text(s, "MedRecord", 1, 1.2, 11, 1.2, font_size=54, bold=True, color=WHITE)
add_text(s, "India's Personal Health Records Platform", 1, 2.5, 11, 0.7,
         font_size=24, color=RGBColor(0xA8, 0xD8, 0xFF))
add_text(s, "Built for 1.4 billion people  |  Powered by ABHA  |  Bridging India's health data gap",
         1, 3.3, 11, 0.6, font_size=16, color=RGBColor(0xCC, 0xE8, 0xFF), italic=True)
add_text(s, "Investor Presentation  —  2026", 1, 6.1, 11, 0.6,
         font_size=14, color=WHITE)

# ── SLIDE 2 — Problem ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_rect(s, 0, 1.1, 13.33, 6.4, LIGHT_BG)
add_text(s, "The Problem", 0.5, 0.15, 12, 0.8, font_size=30, bold=True, color=WHITE)

problems = [
    ("70%", "of India's healthcare is delivered by unorganised clinics & chemists\nnot connected to any digital system"),
    ("85%", "of patients carry physical paper records — prescriptions, lab reports,\ndischarge summaries — prone to loss and damage"),
    ("0%", "of records from small/local providers appear automatically in\nABHA / ABDM — the government's own health data network"),
    ("300M+", "middle-class Indians are NOT primary beneficiaries of PM-JAY,\nyet manage the most complex personal health record portfolios"),
]

colors_box = [TEAL, NAVY, RGBColor(0xC0, 0x39, 0x2B), AMBER]
for i, (stat, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    lx = 0.5 + col * 6.5
    ty = 1.35 + row * 2.6
    add_rect(s, lx, ty, 6.1, 2.3, colors_box[i])
    add_text(s, stat, lx+0.2, ty+0.1, 5.7, 0.8, font_size=40, bold=True, color=WHITE)
    add_text(s, desc, lx+0.2, ty+0.9, 5.6, 1.2, font_size=13, color=WHITE, wrap=True)

# ── SLIDE 3 — Market Opportunity ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 1.1, TEAL)
add_rect(s, 0, 1.1, 13.33, 6.4, LIGHT_BG)
add_text(s, "Market Opportunity", 0.5, 0.15, 12, 0.8, font_size=30, bold=True, color=WHITE)

stats = [
    ("$372B", "Indian Healthcare\nMarket by 2030"),
    ("300M+", "Middle-class Indians\nnot covered by PM-JAY"),
    ("560M", "Smartphone users\nin India (2026)"),
    ("ABHA", "500M+ Health IDs\nissued — yet no\nconsumer UX layer"),
]
for i, (num, label) in enumerate(stats):
    lx = 0.5 + i * 3.2
    add_rect(s, lx, 1.5, 2.9, 2.4, NAVY)
    add_text(s, num, lx+0.15, 1.65, 2.6, 0.9, font_size=30, bold=True, color=TEAL)
    add_text(s, label, lx+0.15, 2.5, 2.6, 0.9, font_size=12, color=WHITE, wrap=True)

add_text(s, "Target Segments", 0.5, 4.15, 12, 0.5, font_size=18, bold=True, color=NAVY)
segs = [
    "Middle-class families managing records for parents, children & self",
    "Chronic disease patients (diabetes, hypertension, cardiac) — 101M diabetics alone",
    "Corporates / HR teams managing employee health documentation",
    "Health-conscious urban millennials seeking digital health organisation",
]
bullet_box(s, segs, 0.5, 4.6, 12.3, 2.5, font_size=14, color=DARK_TEXT)

# ── SLIDE 4 — The ABHA Gap ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_rect(s, 0, 1.1, 13.33, 6.4, LIGHT_BG)
add_text(s, "The ABHA Gap — Our Core Opportunity", 0.5, 0.15, 12.5, 0.8,
         font_size=28, bold=True, color=WHITE)

# Left column — what ABHA does
add_rect(s, 0.4, 1.3, 5.8, 5.6, NAVY)
add_text(s, "What ABHA Provides", 0.6, 1.4, 5.5, 0.55,
         font_size=16, bold=True, color=TEAL)
abha_has = [
    "14-digit national health ID for every citizen",
    "Links records from registered government hospitals",
    "CGHS, ESIC, AIIMS — connected as HIPs",
    "Open ABDM API ecosystem",
    "500M+ IDs issued (growing)",
]
bullet_box(s, abha_has, 0.6, 2.0, 5.4, 4.0, font_size=13, color=WHITE)

# Right column — the gap
add_rect(s, 6.8, 1.3, 6.1, 5.6, RGBColor(0xC0, 0x39, 0x2B))
add_text(s, "What ABHA Misses (Our Space)", 7.0, 1.4, 5.8, 0.55,
         font_size=16, bold=True, color=WHITE)
gap_items = [
    "Records from 70%+ unregistered local clinics & chemists",
    "Manual / paper prescriptions & lab reports",
    "Family-level record management (ABHA is individual-only)",
    "Consumer-grade UX — ABHA portal is bureaucratic",
    "Offline access in low-connectivity Tier 2/3 cities",
    "Insurance claims bridge — linking records to insurers",
]
bullet_box(s, gap_items, 7.0, 2.0, 5.7, 4.5, font_size=13, color=WHITE)

# Center arrow
add_text(s, "GAP\n=\nOUR\nMARKET", 6.0, 2.8, 0.8, 2.5,
         font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── SLIDE 5 — Solution ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 1.1, TEAL)
add_rect(s, 0, 1.1, 13.33, 6.4, LIGHT_BG)
add_text(s, "Our Solution — MedRecord", 0.5, 0.15, 12, 0.8,
         font_size=30, bold=True, color=WHITE)

add_text(s, "ABHA is India's health data backbone. MedRecord is the consumer layer on top.",
         0.5, 1.25, 12.3, 0.55, font_size=16, italic=True, bold=True, color=NAVY)

features = [
    (NAVY, "ABHA Integration", [
        "Link records to your 14-digit ABHA ID",
        "Pull records from ABDM-registered hospitals",
        "Push captured records into national ABHA ecosystem",
        "Become ABDM-certified HIU / HIP",
    ]),
    (TEAL, "Capture What ABHA Misses", [
        "Camera scan with OCR text extraction (ML Kit)",
        "Gallery upload — JPG, PNG",
        "PDF file upload and preview",
        "Manual text entry for verbal prescriptions",
    ]),
    (GREEN, "Organise & Access", [
        "Family record management — one login, all members",
        "Categorise: Prescriptions / Lab Reports / Discharge",
        "Full-text search across all records",
        "Offline access — local storage, always available",
    ]),
]

for i, (color, title, bullets) in enumerate(features):
    lx = 0.4 + i * 4.3
    add_rect(s, lx, 2.0, 4.0, 5.0, color)
    add_text(s, title, lx+0.15, 2.1, 3.7, 0.55,
             font_size=15, bold=True, color=WHITE)
    bullet_box(s, bullets, lx+0.15, 2.75, 3.7, 4.0, font_size=12, color=WHITE)

# ── SLIDE 6 — ABHA Differentiator (Deep Dive) ────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_rect(s, 0, 1.1, 13.33, 6.4, LIGHT_BG)
add_text(s, "ABHA Differentiator — Why This Is Our Moat", 0.5, 0.15, 12.5, 0.8,
         font_size=27, bold=True, color=WHITE)

# The strategic insight box
add_rect(s, 0.4, 1.2, 12.5, 1.0, TEAL)
add_text(s,
    "The government builds infrastructure. Private companies build experiences on top.  "
    "UPI -> PhonePe/Paytm  |  DigiLocker -> Apps  |  ABHA -> MedRecord",
    0.6, 1.3, 12.1, 0.7, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

moat_points = [
    (AMBER, "ABDM HIU Certification",
     "Legally pull records from the ENTIRE ABDM network — government hospitals, CGHS, ESIC, registered private chains. "
     "No competitor has done this at consumer scale."),
    (RGBColor(0x8E, 0x44, 0xAD), "Capture the Informal 70%",
     "70% of India's healthcare is informal and invisible to ABHA. "
     "Our OCR-powered capture is the ONLY bridge between paper India and digital India."),
    (GREEN, "Family Account Layer",
     "ABHA is strictly one-ID-one-person. Families need a single pane of glass. "
     "No government solution addresses this. We own this use case."),
    (RGBColor(0xC0, 0x39, 0x2B), "Insurance & Fintech Bridge",
     "Health insurers pay for fast, accurate claims. Linking records to ABHA-verified identity "
     "reduces fraud and claim processing time — a direct revenue opportunity."),
]

for i, (color, title, desc) in enumerate(moat_points):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.45
    ty = 2.4 + row * 2.35
    add_rect(s, lx, ty, 0.12, 1.9, color)
    add_text(s, title, lx+0.25, ty+0.05, 5.9, 0.45, font_size=14, bold=True, color=NAVY)
    add_text(s, desc, lx+0.25, ty+0.5, 5.9, 1.35, font_size=12, color=DARK_TEXT, wrap=True)

# ── SLIDE 7 — Business Model ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 1.1, TEAL)
add_rect(s, 0, 1.1, 13.33, 6.4, LIGHT_BG)
add_text(s, "Business Model", 0.5, 0.15, 12, 0.8, font_size=30, bold=True, color=WHITE)

streams = [
    (NAVY,  "B2C Freemium",        "Free: 50 records, 1 family member\nPro: ₹299/mo — Unlimited + family\nFamily: ₹499/mo — up to 6 members", "Year 1 focus"),
    (TEAL,  "B2B Corporate",       "HR teams & corporates — employee health\nrecord management & wellness programmes\n₹150/employee/year", "Year 2"),
    (GREEN, "Insurance Partners",  "Revenue share with health insurers\nfor claims acceleration & fraud reduction\nusing ABHA-verified record bundles", "Year 2"),
    (AMBER, "Diagnostic Labs",     "Labs pay to receive digitised records\nand appear as preferred scan destinations\nin the app", "Year 3"),
    (RGBColor(0x8E, 0x44, 0xAD), "Telemedicine",
     "In-app doctor consultations with\nfull record context. Commission per\nconsultation or subscription bundle", "Year 3"),
]

for i, (color, title, desc, timing) in enumerate(streams):
    lx = 0.3 + i * 2.55
    add_rect(s, lx, 1.3, 2.35, 5.7, color)
    add_text(s, title, lx+0.1, 1.4, 2.15, 0.5, font_size=13, bold=True, color=WHITE)
    add_text(s, timing, lx+0.1, 1.9, 2.15, 0.35,
             font_size=10, italic=True, color=RGBColor(0xCC, 0xFF, 0xEE))
    add_text(s, desc, lx+0.1, 2.3, 2.15, 4.5, font_size=11, color=WHITE, wrap=True)

# ── SLIDE 8 — Competitive Landscape ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_rect(s, 0, 1.1, 13.33, 6.4, LIGHT_BG)
add_text(s, "Competitive Landscape", 0.5, 0.15, 12, 0.8, font_size=30, bold=True, color=WHITE)

headers = ["", "ABHA\n(Govt)", "HealthifyMe", "Practo\nHealth Vault", "MedRecord\n(Us)"]
rows = [
    ["ABHA ID Integration",  "Native",  "No",    "Partial",  "Deep + HIU"],
    ["Captures Informal Records", "No", "No",    "No",       "Yes (OCR)"],
    ["Family Management",    "No",       "No",    "No",       "Yes"],
    ["Offline Access",       "No",       "No",    "No",       "Yes"],
    ["Insurance Bridge",     "No",       "No",    "No",       "Yes (planned)"],
    ["B2B / Corporate",      "No",       "No",    "No",       "Yes (Year 2)"],
    ["Consumer UX",          "Basic",    "Good",  "Average",  "Premium"],
]

col_colors = [LIGHT_BG, RGBColor(0xDD,0xDD,0xDD), RGBColor(0xDD,0xDD,0xDD),
              RGBColor(0xDD,0xDD,0xDD), TEAL]
col_widths = [3.5, 1.8, 1.8, 2.0, 2.3]
col_starts = [0.3, 3.85, 5.65, 7.45, 9.65]

for ci, (hdr, cw, cx) in enumerate(zip(headers, col_widths, col_starts)):
    add_rect(s, cx, 1.2, cw-0.05, 0.65,
             NAVY if ci == 0 else (TEAL if ci == 4 else RGBColor(0x55,0x55,0x55)))
    fc = WHITE
    add_text(s, hdr, cx+0.05, 1.22, cw-0.1, 0.6,
             font_size=11, bold=True, color=fc, align=PP_ALIGN.CENTER, wrap=True)

for ri, row in enumerate(rows):
    row_bg = WHITE if ri % 2 == 0 else RGBColor(0xEE, 0xF4, 0xFF)
    ty = 1.9 + ri * 0.63
    for ci, (cell, cw, cx) in enumerate(zip(row, col_widths, col_starts)):
        bg = row_bg if ci < 4 else RGBColor(0xE0, 0xF7, 0xF5)
        add_rect(s, cx, ty, cw-0.05, 0.6, bg)
        cell_color = DARK_TEXT
        if ci == 4:
            cell_color = GREEN if cell not in ["No", "Basic", "Average"] else RGBColor(0xC0, 0x39, 0x2B)
        elif cell == "No":
            cell_color = RGBColor(0xAA, 0xAA, 0xAA)
        add_text(s, cell, cx+0.07, ty+0.08, cw-0.15, 0.48,
                 font_size=11, color=cell_color,
                 bold=(ci == 4), align=PP_ALIGN.CENTER)

# ── SLIDE 9 — Traction & Roadmap ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 1.1, TEAL)
add_rect(s, 0, 1.1, 13.33, 6.4, LIGHT_BG)
add_text(s, "Roadmap", 0.5, 0.15, 12, 0.8, font_size=30, bold=True, color=WHITE)

phases = [
    (NAVY, "Phase 1  Q1-Q2 2026", [
        "MVP launch — Android + iOS",
        "Core record capture + OCR",
        "ABHA ID linking (manual)",
        "Firebase cloud sync",
        "Target: 10,000 downloads",
    ]),
    (TEAL, "Phase 2  Q3-Q4 2026", [
        "Apply for ABDM HIU certification",
        "Pull records from ABHA ecosystem",
        "Family account management",
        "Health insurance partner integrations",
        "Target: 100,000 MAU",
    ]),
    (GREEN, "Phase 3  2027", [
        "ABDM HIP certification",
        "B2B corporate health product",
        "Telemedicine in-app consultations",
        "Diagnostic lab partnerships",
        "Target: 1M MAU  |  ₹5Cr ARR",
    ]),
]

for i, (color, phase, items) in enumerate(phases):
    lx = 0.5 + i * 4.25
    add_rect(s, lx, 1.3, 4.0, 5.8, color)
    add_text(s, phase, lx+0.15, 1.4, 3.7, 0.55, font_size=14, bold=True, color=WHITE)
    bullet_box(s, items, lx+0.15, 2.05, 3.7, 4.8, font_size=13, color=WHITE)

# ── SLIDE 10 — Ask ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 5.9, 13.33, 1.6, TEAL)

add_text(s, "The Ask", 1, 1.0, 11, 1.0, font_size=48, bold=True, color=WHITE)
add_text(s, "Seed Round  —  ₹2.5 Crore  (~ $300K USD)", 1, 2.1, 11, 0.65,
         font_size=26, color=TEAL, bold=True)

use_of_funds = [
    "40%  —  Product & Engineering (ABDM HIU certification, iOS, offline sync)",
    "25%  —  Growth & User Acquisition (Tier 1 & Tier 2 cities)",
    "20%  —  Partnerships (Insurance, Diagnostic Labs, Corporate)",
    "15%  —  Operations & Compliance (DPDP Act, healthcare regulations)",
]
bullet_box(s, use_of_funds, 1, 2.9, 11.3, 2.6, font_size=15, color=WHITE)

add_text(s, "We are building the consumer experience layer on India's national health infrastructure.",
         1, 6.1, 11.3, 0.65, font_size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/maximus/projects/llm-lab/MedRecord_Investor_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
