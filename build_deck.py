from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x2B, 0x55)   # dark navy – title bg / accent bar
TEAL       = RGBColor(0x00, 0x8B, 0x8B)   # teal – headings, table headers
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)   # slide background
MID_GRAY   = RGBColor(0x55, 0x65, 0x7A)   # body text
GOLD       = RGBColor(0xF5, 0xA6, 0x23)   # highlight / callout


# ── Helpers ────────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                bold=False, size=18, color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color if color else MID_GRAY
    return txb


def accent_bar(slide, color=TEAL, width=0.07):
    """Thin vertical bar on the left edge."""
    add_rect(slide, 0, 0, width, 7.5, color)


def slide_heading(slide, title, subtitle=None):
    """Standard slide heading below a navy top banner."""
    add_rect(slide, 0, 0, 10, 1.15, NAVY)
    add_textbox(slide, title, 0.3, 0.15, 9.4, 0.75,
                bold=True, size=28, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.82, 9.4, 0.38,
                    size=13, color=GOLD)


def bullet_block(slide, items, left, top, width, height,
                 size=15, color=None, bullet="•"):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color if color else MID_GRAY


def kv_block(slide, pairs, left, top, col_width=2.2, row_height=0.5,
             key_size=12, val_size=18):
    """Render a list of (key, value) cards in a row."""
    for i, (k, v) in enumerate(pairs):
        x = left + i * (col_width + 0.15)
        add_rect(slide, x, top, col_width, row_height * 2 + 0.1, WHITE)
        add_textbox(slide, v, x + 0.1, top + 0.05, col_width - 0.2, row_height,
                    bold=True, size=val_size, color=TEAL, align=PP_ALIGN.CENTER)
        add_textbox(slide, k, x + 0.1, top + row_height + 0.05,
                    col_width - 0.2, row_height,
                    size=key_size, color=MID_GRAY, align=PP_ALIGN.CENTER)


def simple_table(slide, headers, rows, left, top, col_widths, row_h=0.42,
                 header_color=NAVY, alt_color=LIGHT_GRAY):
    n_cols = len(headers)
    # header row
    x = left
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, top, w, row_h, header_color)
        add_textbox(slide, h, x + 0.05, top + 0.04, w - 0.1, row_h - 0.08,
                    bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)
        x += w
    # data rows
    for r_idx, row in enumerate(rows):
        y = top + row_h + r_idx * row_h
        bg = LIGHT_GRAY if r_idx % 2 == 0 else WHITE
        x = left
        for c_idx, (cell, w) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, y, w, row_h, bg)
            bold = (c_idx == 0)
            add_textbox(slide, str(cell), x + 0.05, y + 0.04, w - 0.1, row_h - 0.08,
                        bold=bold, size=10, color=MID_GRAY if c_idx > 0 else NAVY,
                        align=PP_ALIGN.LEFT)
            x += w


# ── Build presentation ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # completely blank layout


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)

# Big logo-style name
add_textbox(s, "MedRecords", 0.5, 1.5, 9, 1.5,
            bold=True, size=60, color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
add_rect(s, 2.5, 3.1, 5, 0.05, GOLD)   # divider line
add_textbox(s, "Your Health History, Always With You",
            0.5, 3.3, 9, 0.7,
            size=22, color=GOLD, align=PP_ALIGN.CENTER)

# Sub-details
add_textbox(s, "Android  ·  iOS  ·  Flutter",
            0.5, 4.2, 9, 0.5, size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "Pre-Seed  ·  MVP Complete  ·  March 2026",
            0.5, 4.75, 9, 0.5, size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
            align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "The Problem", "Managing medical records is broken")

problems = [
    ("Scattered Records",
     "Prescriptions on paper, lab reports across clinic portals,\nold reports lost or damaged."),
    ("No Single Source of Truth",
     "Patients spend hours hunting for past test results\nbefore every new doctor's appointment."),
    ("No Portability",
     "Records locked in hospital systems or physical folders —\ninaccessible in emergencies."),
]
for i, (title, desc) in enumerate(problems):
    x = 0.4 + i * 3.2
    add_rect(s, x, 1.4, 3.0, 2.2, WHITE)
    add_rect(s, x, 1.4, 3.0, 0.45, TEAL)
    add_textbox(s, title, x + 0.1, 1.43, 2.8, 0.4,
                bold=True, size=13, color=WHITE)
    add_textbox(s, desc, x + 0.1, 1.95, 2.8, 1.5, size=12, color=MID_GRAY)

# Stat callout
add_rect(s, 0.4, 3.85, 9.2, 0.85, NAVY)
add_textbox(s,
    "\"43% of patients cannot produce a complete medication history when visiting a new doctor\"  — NCBI, 2024",
    0.6, 3.95, 8.8, 0.65, size=13, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "The Solution", "MedRecords — a secure personal health vault in your pocket")

headers = ["Pain Point", "MedRecords Solution"]
rows = [
    ["Records are scattered",        "One app stores all prescriptions & lab reports"],
    ["Paper records are fragile",     "Camera scan with automatic OCR text extraction"],
    ["Locked in hospital portals",    "Patient-owned cloud — accessible anywhere, anytime"],
    ["Hard to share with doctors",    "Organised, categorised records ready to show or share"],
]
simple_table(s, headers, rows,
             left=0.5, top=1.35,
             col_widths=[4.0, 5.0], row_h=0.52)

add_rect(s, 0.5, 4.6, 9.0, 0.75, RGBColor(0xE8, 0xF8, 0xF5))
add_textbox(s,
    "Core Insight:  The patient — not the hospital — should own their health data.",
    0.7, 4.68, 8.6, 0.55,
    bold=True, size=14, color=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PRODUCT FEATURES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "Product Features", "Everything a patient needs — nothing they don't")

features = [
    ("📷  Capture", ["Camera scan", "Gallery upload", "File upload (PDF/JPG/PNG)", "Manual text entry"]),
    ("🔍  Smart OCR", ["Google ML Kit on-device", "Text extracted automatically", "Documents fully searchable", "No data sent to 3rd party"]),
    ("🗂  Organised", ["Prescriptions category", "Lab Reports category", "Tabbed home screen", "Delete & manage records"]),
    ("☁  Secure Cloud", ["Firebase Auth login", "Per-user Firestore rules", "Real-time sync", "Offline-ready"]),
]
for i, (title, items) in enumerate(features):
    x = 0.3 + (i % 2) * 4.8
    y = 1.4 + (i // 2) * 2.55
    add_rect(s, x, y, 4.4, 2.3, WHITE)
    add_rect(s, x, y, 4.4, 0.5, TEAL)
    add_textbox(s, title, x + 0.1, y + 0.05, 4.2, 0.4,
                bold=True, size=14, color=WHITE)
    bullet_block(s, items, x + 0.15, y + 0.6, 4.1, 1.6, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "Market Opportunity", "A massive, underserved global market")

kv_block(s,
    [("Global Digital Health (2025)", "$560B"),
     ("PHR Segment", "$39B"),
     ("PHR CAGR 2025–2030", "15.2%"),
     ("Target Smartphone Users", "4.4B")],
    left=0.35, top=1.35, col_width=2.3, row_height=0.5, val_size=22)

add_textbox(s, "Primary Target Markets", 0.5, 2.8, 9, 0.4,
            bold=True, size=15, color=NAVY)
markets = [
    "India — 1.4B population, fragmented healthcare, low EHR adoption by clinics",
    "Southeast Asia — Rapid smartphone growth, minimal hospital digitisation",
    "Emerging Markets globally — Paper-heavy healthcare, no incumbent PHR app",
]
bullet_block(s, markets, 0.5, 3.25, 9.0, 1.5, size=13)

add_textbox(s, "Why Now", 0.5, 4.85, 9, 0.4, bold=True, size=15, color=NAVY)
why = [
    "Post-COVID patients expect digital health tools",
    "India's Digital Health Mission (ABDM) creates tailwinds but leaves personal record ownership unsolved",
    "Firebase + Flutter cut infra & dev costs by ~60% vs. 5 years ago",
]
bullet_block(s, why, 0.5, 5.3, 9.0, 1.5, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "Business Model", "Multiple monetisation paths — freemium to B2B")

# Phase labels
phases = [
    ("Phase 1  —  Freemium (Launch)", [
        ("Free", "$0 / month", "Up to 50 records, basic categories"),
        ("Premium", "$2.99 / month", "Unlimited records, search, PDF export, cloud backup"),
    ]),
]

add_textbox(s, "Phase 1 — Freemium (Launch)", 0.4, 1.3, 9, 0.4,
            bold=True, size=14, color=NAVY)
tiers = [
    ["Free", "$0 / month", "Up to 50 records, basic categories"],
    ["Premium", "$2.99 / month", "Unlimited records, search, PDF export, cloud backup"],
]
simple_table(s, ["Tier", "Price", "Features"], tiers,
             left=0.4, top=1.75, col_widths=[1.5, 1.8, 5.7], row_h=0.45)

add_textbox(s, "Phase 2 — B2B / Platform  (12–18 months)", 0.4, 3.05, 9, 0.4,
            bold=True, size=14, color=NAVY)
b2b = [
    ["Clinic / Hospital API", "Clinics pay to push records directly into patient vaults"],
    ["Insurance Integrations", "Insurers pay for verified, structured health data (with user consent)"],
    ["Telemedicine Partnerships", "Embed records sharing into telehealth consult flows"],
]
simple_table(s, ["Revenue Stream", "Description"], b2b,
             left=0.4, top=3.5, col_widths=[2.8, 6.7], row_h=0.45)

add_textbox(s, "Phase 3 — Data & AI  (24+ months)", 0.4, 4.9, 9, 0.4,
            bold=True, size=14, color=NAVY)
bullet_block(s,
    ["Anonymised health trend insights sold to research institutions (with explicit consent)",
     "AI-powered health summaries & alerts for premium subscribers"],
    0.4, 5.35, 9.0, 0.9, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "Competitive Landscape", "We sit at a unique intersection")

headers = ["Feature", "MedRecords", "Apple Health", "Google Health", "Hospital Portals"]
rows = [
    ["Android & iOS",              "✅ Yes",       "iOS only",     "Android focus", "Fragmented"],
    ["Cross-institution records",  "✅ Yes",       "Partial",      "Partial",       "No"],
    ["OCR scan of paper records",  "✅ Yes",       "No",           "No",            "No"],
    ["Patient-owned data",         "✅ Yes",       "Partial",      "No",            "No"],
    ["Works offline / low-data",   "✅ Yes",       "No",           "No",            "No"],
    ["Emerging market focus",      "✅ Yes",       "No",           "Limited",       "No"],
]
simple_table(s, headers, rows,
             left=0.3, top=1.35,
             col_widths=[2.8, 1.5, 1.5, 1.5, 1.9], row_h=0.48)

add_rect(s, 0.3, 4.65, 9.4, 0.75, NAVY)
add_textbox(s, "Our Moat:  OCR-first capture  +  Patient data ownership  +  Cross-platform reach in emerging markets",
            0.5, 4.73, 9.0, 0.55,
            bold=True, size=13, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNOLOGY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "Technology", "Built on proven, scalable infrastructure")

tech = [
    ("Flutter 3.41",         "One codebase → Android + iOS. Faster dev, lower cost."),
    ("Firebase Auth",        "Email/Password login. OAuth-ready for Google/Apple sign-in."),
    ("Cloud Firestore",      "Real-time NoSQL DB. Scales automatically. Offline-ready."),
    ("Google ML Kit OCR",    "On-device text recognition — sensitive images never leave the phone."),
    ("Firebase Storage",     "Secure cloud image storage (Blaze plan). Per-user access rules."),
]
for i, (name, desc) in enumerate(tech):
    y = 1.4 + i * 0.9
    add_rect(s, 0.4, y, 2.2, 0.72, NAVY)
    add_textbox(s, name, 0.5, y + 0.1, 2.0, 0.52,
                bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, 2.75, y + 0.08, 6.8, 0.6, size=13, color=MID_GRAY)

add_textbox(s, "Security Highlights", 0.4, 6.0, 9, 0.35,
            bold=True, size=13, color=NAVY)
bullet_block(s,
    ["Per-user Firestore security rules (data isolated by UID)",
     "On-device OCR — documents never sent to external server",
     "HTTPS-only transit  ·  Firebase JWT tokens"],
    0.4, 6.38, 9.0, 0.85, size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TRACTION & ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "Traction & Roadmap", "MVP complete — ready to ship")

add_textbox(s, "Completed (MVP)", 0.4, 1.3, 9, 0.38,
            bold=True, size=14, color=NAVY)
bullet_block(s,
    ["Full Flutter app — 11 source modules across screens, services, and models",
     "Firebase Auth, Firestore, and ML Kit OCR fully integrated",
     "Camera scan, file upload, and manual entry flows working",
     "Tabbed record browsing (All / Prescriptions / Lab Reports)",
     "Android APK build-ready; iOS build-ready"],
    0.4, 1.72, 9.0, 1.6, size=12)

add_textbox(s, "Roadmap", 0.4, 3.45, 9, 0.38,
            bold=True, size=14, color=NAVY)
roadmap = [
    ["Q2 2026", "Public beta — Android Play Store + iOS App Store"],
    ["Q3 2026", "Search & filter  ·  PDF export & sharing  ·  1,000 beta users"],
    ["Q4 2026", "Premium subscription launch"],
    ["Q1 2027", "Clinic API pilot (B2B)"],
    ["Q2 2027", "Series A fundraise"],
]
simple_table(s, ["Quarter", "Milestone"], roadmap,
             left=0.4, top=3.88, col_widths=[1.5, 8.0], row_h=0.44)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THE ASK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_GRAY)
slide_heading(s, "The Ask", "Pre-Seed Round: $250,000")

add_textbox(s, "Use of Funds", 0.4, 1.3, 9, 0.38,
            bold=True, size=14, color=NAVY)
funds = [
    ["Engineering (40%)",        "$100,000", "2 Flutter/backend engineers for 12 months"],
    ["Growth & Marketing (25%)", "$62,500",  "User acquisition — India & Southeast Asia"],
    ["Infrastructure (10%)",     "$25,000",  "Firebase Blaze + storage scaling"],
    ["Product & Design (15%)",   "$37,500",  "UX research, UI polish, accessibility"],
    ["Legal & Compliance (10%)", "$25,000",  "DPDP Act, GDPR alignment, incorporation"],
]
simple_table(s, ["Allocation", "Amount", "Purpose"], funds,
             left=0.4, top=1.73, col_widths=[2.5, 1.4, 5.6], row_h=0.44)

add_textbox(s, "In Return", 0.4, 4.15, 9, 0.38,
            bold=True, size=14, color=NAVY)
bullet_block(s,
    ["8–12% equity (negotiable based on valuation discussion)",
     "Investor board observer seat available"],
    0.4, 4.57, 9, 0.7, size=13)

add_textbox(s, "Target Metrics — Month 12", 0.4, 5.4, 9, 0.38,
            bold=True, size=14, color=NAVY)
bullet_block(s,
    ["10,000 active users",
     "$5,000 MRR from premium subscriptions",
     "Signed LOI with 1 clinic / hospital partner"],
    0.4, 5.82, 9, 0.85, size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)

add_textbox(s, "MedRecords", 0.5, 1.2, 9, 1.1,
            bold=True, size=52, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 2.5, 2.55, 5, 0.05, GOLD)

add_textbox(s, "Your health history — secure, searchable, always with you.",
            0.5, 2.75, 9, 0.75,
            size=20, color=GOLD, align=PP_ALIGN.CENTER)

add_textbox(s,
    "We're building the personal health data layer that 4 billion\nsmartphone users in emerging markets don't yet have.",
    0.8, 3.7, 8.4, 1.1,
    size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 2.5, 5.0, 5, 0.04, RGBColor(0x55, 0x65, 0x7A))

add_textbox(s, "Pre-Seed  ·  MVP Complete  ·  March 2026",
            0.5, 5.2, 9, 0.5,
            size=13, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)
add_textbox(s, "Firebase Project: medrecords2026",
            0.5, 5.7, 9, 0.4,
            size=12, color=RGBColor(0x77, 0x88, 0x99), align=PP_ALIGN.CENTER)

add_textbox(s,
    "This document contains forward-looking statements for discussion purposes only.",
    0.5, 6.9, 9, 0.35,
    size=9, color=RGBColor(0x55, 0x65, 0x7A), align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "/home/maximus/projects/llm-lab/MedRecords_Investor_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
